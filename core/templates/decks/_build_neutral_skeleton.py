"""Generate ``neutral_skeleton.pptx`` — brand-agnostic deck template.

Gives the default ``python-pptx`` slide layouts + placeholders convenient
generic names (``Title``, ``TitleAndContent`` …) so simple example specs
can address them by name. The engine requires NO specific names — it
discovers a template's real layouts (see ``core/conventions/deck-template.md``
and ``deck-catalog``); these names are just a convenience for this
brand-agnostic skeleton. The ``.pptx`` artifact is committed alongside this
script; rerun the script if you change the mapping.

Usage::

    python core/templates/decks/_build_neutral_skeleton.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


OUTPUT = Path(__file__).resolve().parent / "neutral_skeleton.pptx"

# Default python-pptx layouts mapped to the ade-ops convention names.
# Layout idx 6 (Blank) is kept as-is and ignored by the engine.
LAYOUT_RENAME = {
    0: "Title",            # Title Slide          (title + subtitle)
    1: "TitleAndContent",  # Title and Content    (title + body bullets)
    2: "Section",          # Section Header       (title divider)
    3: "TwoContent",       # Two Content          (title + two body cols)
    4: "Quote",            # Comparison           (repurposed for quote+attribution)
    5: "Closing",          # Title Only           (closing / contact slide)
    7: "ImageRight",       # Content with Caption (body + small picture)
    8: "ImageFull",        # Picture with Caption (large picture + caption)
}

# Per-layout placeholder rename map: (layout idx) -> {idx: convention-name}.
# The engine matches placeholders by name first, idx-fallback second. Naming
# them per the convention removes the fallback warnings on freshly built
# decks. Idx mapping reflects the default python-pptx layout shapes.
PLACEHOLDER_RENAME: dict[int, dict[int, str]] = {
    0: {0: "title", 1: "subtitle"},          # Title
    1: {0: "title", 1: "body"},               # TitleAndContent
    2: {0: "title", 1: "subtitle"},           # Section (Section Header has subtitle slot)
    3: {0: "title", 1: "left", 2: "right"},   # TwoContent
    4: {0: "title", 1: "quote", 3: "attribution"},  # Quote (Comparison: 2 headers + 2 bodies; use 1 + 3)
    5: {0: "title"},                          # Closing (Title Only — no body)
    7: {0: "title", 1: "body", 2: "image"},   # ImageRight (Content with Caption)
    8: {0: "image", 1: "caption"},            # ImageFull (Picture with Caption)
}


def main() -> None:
    prs = Presentation()
    layouts = prs.slide_layouts
    for layout_idx, new_layout_name in LAYOUT_RENAME.items():
        if layout_idx >= len(layouts):
            continue
        layout = layouts[layout_idx]
        layout.name = new_layout_name
        for ph in layout.placeholders:
            new_ph_name = PLACEHOLDER_RENAME.get(layout_idx, {}).get(
                ph.placeholder_format.idx
            )
            if new_ph_name:
                ph.name = new_ph_name
    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT}")
    for layout in prs.slide_layouts:
        ph_names = [ph.name for ph in layout.placeholders]
        print(f"  {layout.name:24} placeholders={ph_names}")


if __name__ == "__main__":
    main()
