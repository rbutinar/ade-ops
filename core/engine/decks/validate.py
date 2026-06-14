"""Template discovery + spec validation for the decks engine.

Two operations, both **non-destructive** and built around *real-layout
discovery* (TICK-035) — there is no required-layout convention:

- :func:`catalog_template` — read a template and list every layout
  (index, name) and its placeholders (idx, name, type). This is what an
  operator or agent reads to author a spec; it replaces the old
  "rename your layouts to 8 fixed names in the Slide Master" step.
- :func:`validate_spec` — check that every layout (by index or name) and
  every placeholder idx a spec references actually exists in its template,
  so missing references surface before the build instead of silently
  dropping content.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from .config import DecksConfig, load_decks_config


@dataclass
class PlaceholderInfo:
    """One placeholder on a layout."""

    idx: int
    name: str
    type: str


@dataclass
class LayoutInfo:
    """One slide layout in a template, with its placeholders."""

    index: int
    name: str
    placeholders: list[PlaceholderInfo] = field(default_factory=list)


def _placeholder_type_name(placeholder_format: Any) -> str:
    """Human-readable placeholder type (e.g. ``TITLE``, ``BODY``,
    ``PICTURE``), tolerant of ``None``."""
    ptype = placeholder_format.type
    if ptype is None:
        return "?"
    return getattr(ptype, "name", str(ptype))


def catalog_template(template_path: Path) -> list[LayoutInfo]:
    """Open ``template_path`` and return one :class:`LayoutInfo` per slide
    layout, in template order (so the index is the value a spec uses in
    ``layout: <index>``)."""
    # Lazy import so the rest of the engine imports without python-pptx.
    from pptx import Presentation

    prs = Presentation(str(template_path))
    catalog: list[LayoutInfo] = []
    for index, layout in enumerate(prs.slide_layouts):
        placeholders = [
            PlaceholderInfo(
                idx=ph.placeholder_format.idx,
                name=ph.name or "",
                type=_placeholder_type_name(ph.placeholder_format),
            )
            for ph in layout.placeholders
        ]
        catalog.append(
            LayoutInfo(index=index, name=layout.name, placeholders=placeholders)
        )
    return catalog


@dataclass
class SpecValidationResult:
    """Outcome of :func:`validate_spec`."""

    spec_path: Path
    template_path: Path
    slide_count: int
    issues: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return not self.issues

    def summary(self) -> str:
        if self.ok:
            return (
                f"OK — {self.slide_count} slide(s), every layout + placeholder "
                f"reference resolves against {self.template_path.name}"
            )
        return f"INVALID — {len(self.issues)} unresolved reference(s)"


def _referenced_indices(slide_spec: dict[str, Any]) -> set[int]:
    """Collect every placeholder idx a slide spec references via the
    idx-keyed directives (``text`` / ``bullets`` / ``force_font`` /
    ``image`` with an explicit ``idx``). Named fields are not idx-bound and
    are reported as build-time warnings, not validation errors."""
    indices: set[int] = set()
    for field_name in ("text", "bullets"):
        value = slide_spec.get(field_name)
        if isinstance(value, dict):
            for key in value:
                try:
                    indices.add(int(key))
                except (TypeError, ValueError):
                    pass
    for entry in slide_spec.get("force_font") or []:
        if isinstance(entry, dict) and "idx" in entry:
            try:
                indices.add(int(entry["idx"]))
            except (TypeError, ValueError):
                pass
    image = slide_spec.get("image")
    if isinstance(image, dict) and "idx" in image:
        try:
            indices.add(int(image["idx"]))
        except (TypeError, ValueError):
            pass
    return indices


def validate_spec(
    spec_path: Path, config: DecksConfig | None = None
) -> SpecValidationResult:
    """Resolve a spec's template and check every layout + placeholder idx it
    references exists. Returns a result whose ``issues`` list is empty when
    the spec is buildable against its template."""
    from .builder import load_spec, resolve_spec_template

    spec_path = spec_path.resolve()
    spec = load_spec(spec_path)
    if config is None:
        config = load_decks_config()
    template_path = resolve_spec_template(spec, config)

    catalog = catalog_template(template_path)
    by_index = {info.index: info for info in catalog}
    by_name = {info.name.lower(): info for info in catalog}

    slides_spec = spec.get("slides") or []
    issues: list[str] = []

    for index, slide_spec in enumerate(slides_spec, start=1):
        if not isinstance(slide_spec, dict):
            issues.append(f"slide #{index}: not a mapping")
            continue
        layout_ref = slide_spec.get("layout")
        if layout_ref is None:
            issues.append(f"slide #{index}: missing ``layout``")
            continue

        info: LayoutInfo | None
        if isinstance(layout_ref, bool):
            issues.append(f"slide #{index}: ``layout`` must not be a boolean")
            continue
        if isinstance(layout_ref, int):
            info = by_index.get(layout_ref)
            if info is None:
                issues.append(
                    f"slide #{index}: layout index {layout_ref} out of range "
                    f"(0..{len(catalog) - 1})"
                )
                continue
        else:
            info = by_name.get(str(layout_ref).lower())
            if info is None:
                issues.append(
                    f"slide #{index}: layout '{layout_ref}' not in template"
                )
                continue

        available = {ph.idx for ph in info.placeholders}
        for idx in sorted(_referenced_indices(slide_spec)):
            if idx not in available:
                issues.append(
                    f"slide #{index}: layout '{info.name}' (#{info.index}) has no "
                    f"placeholder idx {idx} (available: "
                    f"{sorted(available) or 'none'})"
                )

    return SpecValidationResult(
        spec_path=spec_path,
        template_path=template_path,
        slide_count=len(slides_spec),
        issues=issues,
    )
