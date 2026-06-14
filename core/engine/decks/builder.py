"""Build a ``.pptx`` from a spec YAML + a brand template.

Convention reference: ``core/conventions/deck-template.md``.

**Real-layout discovery** (TICK-035): the spec addresses the template's
*actual* layouts and placeholders — there is no layout-rename convention.
A slide picks a layout by **index** (``layout: 24``) or by its **real
name** (``layout: "Content: 3 columns"``), and fills placeholders by
**idx** (``text: {0: ...}``) for rich brand templates, or by generic
**named fields** (``title:``, ``body:`` …) for templates whose
placeholders are named. Discover a template's layouts + placeholder
indices with ``deck-catalog``.

Build flow::

    spec.yaml --> DecksConfig.pack_path(pack) --> <pack>/<template>.pptx
                                                       |
                                                       v
                                                  Presentation()
                                                       |
                            clear template example slides (unless kept)
                                                       |
              for each slide: resolve layout (idx|name), add, populate
                                                       |
                                                       v
                                                   PPTX output
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Any

import yaml

from .config import DecksConfig, DecksConfigError, load_decks_config


class BuildError(Exception):
    """Raised when a spec cannot be built into a PPTX (bad layout ref,
    missing template file, unreadable image, ...)."""


# Generic, layout-INDEPENDENT named fields. These are sugar for templates
# whose placeholders carry names (e.g. the shipped neutral skeleton): the
# value is routed to a placeholder matched by *name* (case-insensitive),
# with an idx-based fallback. They are NOT a layout vocabulary — any layout
# may use any field. Rich brand templates address placeholders by idx
# instead (the ``text`` / ``bullets`` / ``image`` / ``force_font`` directives).
# ``kind``: ``text`` (single string), ``body`` (string or list of bullets).
_NAMED_FIELDS: list[tuple[str, tuple[str, ...], str]] = [
    ("title", ("title",), "text"),
    ("subtitle", ("subtitle",), "text"),
    ("body", ("body", "content"), "body"),
    ("left", ("left", "left_content"), "body"),
    ("right", ("right", "right_content"), "body"),
    ("quote", ("quote",), "text"),
    ("attribution", ("attribution",), "text"),
    ("caption", ("caption",), "text"),
]


@dataclass
class BuildReport:
    """Summary of a build operation."""

    spec_path: Path
    template_path: Path
    output_path: Path
    slide_count: int
    warnings: list[str]


def load_spec(spec_path: Path) -> dict[str, Any]:
    """Read + parse a spec YAML, returning its top-level mapping."""
    spec_path = spec_path.resolve()
    if not spec_path.exists():
        raise BuildError(f"Spec file not found: {spec_path}")
    spec = yaml.safe_load(spec_path.read_text(encoding="utf-8")) or {}
    if not isinstance(spec, dict):
        raise BuildError(f"{spec_path}: top-level YAML must be a mapping")
    return spec


def resolve_spec_template(spec: dict[str, Any], config: DecksConfig) -> Path:
    """Resolve the template ``.pptx`` a spec targets via its ``meta``."""
    meta = spec.get("meta") or {}
    pack_name = meta.get("pack") or config.default_pack
    if not pack_name:
        raise BuildError("meta.pack is required (no default_pack in decks.yaml)")
    return _resolve_template_path(config, pack_name, meta.get("template"))


def build_from_spec(
    spec_path: Path,
    output_path: Path | None = None,
    config: DecksConfig | None = None,
) -> BuildReport:
    """Build a PPTX from ``spec_path``.

    If ``output_path`` is not given, the output goes to
    ``DecksConfig.output_dir`` with a name derived from the spec
    basename plus today's date.
    """
    # Lazy import so the engine module is importable without pptx.
    from pptx import Presentation

    spec_path = spec_path.resolve()
    spec = load_spec(spec_path)

    if config is None:
        config = load_decks_config()

    meta = spec.get("meta") or {}
    slides_spec = spec.get("slides") or []
    if not isinstance(slides_spec, list) or not slides_spec:
        raise BuildError(f"{spec_path}: ``slides`` must be a non-empty list")

    template_path = resolve_spec_template(spec, config)

    prs = Presentation(str(template_path))

    # Drop the template's own example slides so the output contains only the
    # spec's slides (masters / layouts / theme are untouched). Brand starter
    # packs ship dozens of example slides; ``keep_template_slides: true`` opts
    # out (e.g. to append onto an existing deck).
    if not meta.get("keep_template_slides", False):
        _clear_slides(prs)

    layouts = list(prs.slide_layouts)
    layouts_by_name: dict[str, Any] = {
        layout.name.lower(): layout for layout in layouts
    }

    warnings: list[str] = []
    for index, slide_spec in enumerate(slides_spec, start=1):
        if not isinstance(slide_spec, dict):
            raise BuildError(f"{spec_path}: slide #{index} must be a mapping")
        layout = _resolve_layout(
            slide_spec, layouts, layouts_by_name, spec_path, index
        )
        slide = prs.slides.add_slide(layout)
        _populate_slide(
            slide,
            slide_spec,
            base_dir=spec_path.parent,
            warnings=warnings,
            slide_index=index,
        )

    if output_path is None:
        output_path = _default_output_path(spec_path, meta, config)
    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    return BuildReport(
        spec_path=spec_path,
        template_path=template_path,
        output_path=output_path,
        slide_count=len(slides_spec),
        warnings=warnings,
    )


# =============================================================================
# Internal helpers
# =============================================================================


def _resolve_template_path(
    config: DecksConfig, pack_name: str, template_name: str | None
) -> Path:
    """Pick the template `.pptx` inside the pack folder.

    If ``template_name`` is ``None`` or ``"default"``, picks the
    alphabetically first ``.pptx`` in the folder.
    """
    try:
        pack_folder = config.pack_path(pack_name)
    except DecksConfigError as exc:
        raise BuildError(str(exc)) from exc

    if not pack_folder.exists():
        raise BuildError(
            f"Pack folder does not exist: {pack_folder}. Check the path in "
            f"{config.config_path} or sync OneDrive."
        )

    if template_name in (None, "default"):
        candidates = sorted(pack_folder.glob("*.pptx"))
        if not candidates:
            raise BuildError(f"No .pptx files in pack folder {pack_folder}")
        return candidates[0]

    candidate = pack_folder / template_name
    if not candidate.exists():
        available = ", ".join(sorted(p.name for p in pack_folder.glob("*.pptx")))
        raise BuildError(
            f"Template '{template_name}' not found in {pack_folder}. "
            f"Available: {available or '(none)'}"
        )
    return candidate


def _default_output_path(
    spec_path: Path, meta: dict[str, Any], config: DecksConfig
) -> Path:
    output_dir = config.output_dir or spec_path.parent
    stamp = date.today().isoformat().replace("-", "")
    return output_dir / f"{spec_path.stem}_{stamp}.pptx"


def _resolve_layout(
    slide_spec: dict[str, Any],
    layouts: list[Any],
    layouts_by_name: dict[str, Any],
    spec_path: Path,
    index: int,
) -> Any:
    """Resolve a slide's ``layout`` reference (int index OR real name)."""
    layout_ref = slide_spec.get("layout")
    if layout_ref is None:
        raise BuildError(
            f"{spec_path}: slide #{index} missing required ``layout`` field"
        )
    if isinstance(layout_ref, bool):  # guard: YAML true/false is not an index
        raise BuildError(
            f"{spec_path}: slide #{index} ``layout`` must be an int index or a "
            f"layout name, got boolean {layout_ref!r}"
        )
    if isinstance(layout_ref, int):
        if not 0 <= layout_ref < len(layouts):
            raise BuildError(
                f"{spec_path}: slide #{index} layout index {layout_ref} out of "
                f"range — template has {len(layouts)} layouts (0..{len(layouts) - 1}). "
                f"Run ``deck-catalog`` to list them."
            )
        return layouts[layout_ref]

    layout = layouts_by_name.get(str(layout_ref).lower())
    if layout is None:
        raise BuildError(
            f"{spec_path}: slide #{index} layout '{layout_ref}' not found in "
            f"template. Available: {', '.join(sorted(layouts_by_name))}. "
            f"Run ``deck-catalog`` for indices + names."
        )
    return layout


def _clear_slides(prs: Any) -> None:
    """Remove the template's example slides (rels + ``sldIdLst`` entries).

    Masters, layouts and theme are untouched — only the populated example
    slides shipped in the template are dropped, so the output contains only
    the spec's slides.
    """
    from pptx.oxml.ns import qn

    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)


def _find_placeholder(slide: Any, names: tuple[str, ...]) -> Any | None:
    """Return the first placeholder whose name matches one of ``names``
    (case-insensitive).

    Looks at the slide first; if no direct hit, falls back to the
    *layout*'s placeholders and maps the match back to the slide by
    ``idx`` (placeholder names live on the layout when the operator
    renames in Slide Master view, and slide placeholders inherit the
    layout's idx but not its name).
    """
    wanted = {n.lower() for n in names}

    for ph in slide.placeholders:
        if (ph.name or "").lower() in wanted:
            return ph

    layout = getattr(slide, "slide_layout", None)
    if layout is not None:
        for layout_ph in layout.placeholders:
            if (layout_ph.name or "").lower() not in wanted:
                continue
            target_idx = layout_ph.placeholder_format.idx
            for slide_ph in slide.placeholders:
                if slide_ph.placeholder_format.idx == target_idx:
                    return slide_ph
    return None


def _find_first_text_placeholder(slide: Any, skip_title: bool = False) -> Any | None:
    """Heuristic fallback for templates whose placeholders are unnamed:
    return the first text-capable placeholder, optionally skipping the
    title (idx 0)."""
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if skip_title and idx == 0:
            continue
        if ph.has_text_frame:
            return ph
    return None


def _find_first_picture_placeholder(slide: Any) -> Any | None:
    """Return the first picture-typed placeholder (PP_PLACEHOLDER.PICTURE
    == 18), or ``None``."""
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 18:
            return ph
    return None


def _set_text(ph: Any, value: str) -> None:
    """Write a single string into a placeholder, replacing any existing text."""
    if ph.has_text_frame:
        ph.text_frame.text = value
    else:
        ph.text = value


def _set_body(ph: Any, value: Any) -> None:
    """Write a body field (string -> single paragraph; list -> bullets)."""
    if not ph.has_text_frame:
        ph.text = str(value)
        return
    tf = ph.text_frame
    tf.clear()
    if isinstance(value, list):
        for i, item in enumerate(value):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(item)
    else:
        tf.text = str(value)


def _set_image(slide: Any, ph: Any, image_path: Path) -> None:
    """Insert an image into a placeholder. If the placeholder is a
    picture placeholder, use ``insert_picture``. Otherwise add a free
    picture at the placeholder's bounding box and remove the placeholder
    (so it does not overlap)."""
    from pptx.util import Emu

    # PP_PLACEHOLDER.PICTURE == 18 — picture-typed placeholder supports
    # ``insert_picture`` directly. Other text placeholders need a manual
    # ``add_picture`` at the same bounding box.
    if ph.placeholder_format.type == 18:  # picture placeholder
        ph.insert_picture(str(image_path))
        return
    left, top = ph.left, ph.top
    width, height = ph.width, ph.height
    sp = ph._element
    sp.getparent().remove(sp)  # remove the now-empty placeholder
    slide.shapes.add_picture(
        str(image_path), Emu(left), Emu(top), width=Emu(width), height=Emu(height)
    )


def _populate_slide(
    slide: Any,
    slide_spec: dict[str, Any],
    base_dir: Path,
    warnings: list[str],
    slide_index: int,
) -> None:
    ph_by_idx = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    # 1. Named-field sugar — route by placeholder name (idx fallback).
    for field_name, placeholder_names, kind in _NAMED_FIELDS:
        if field_name not in slide_spec:
            continue
        value = slide_spec[field_name]
        if value is None:
            continue
        ph = _find_placeholder(slide, placeholder_names)
        if ph is None:
            ph = _find_first_text_placeholder(
                slide, skip_title=(field_name != "title")
            )
            if ph is not None:
                warnings.append(
                    f"slide #{slide_index}: field '{field_name}' fell back to "
                    f"first text placeholder (no placeholder named "
                    f"{'/'.join(placeholder_names)})"
                )
        if ph is None:
            warnings.append(
                f"slide #{slide_index}: field '{field_name}' skipped "
                f"(no matching placeholder)"
            )
            continue
        if kind == "text":
            _set_text(ph, str(value))
        else:  # body
            _set_body(ph, value)

    # 2. idx-keyed text — set placeholder <idx> to a single string.
    for idx, value in _idx_map(slide_spec.get("text"), "text", slide_index).items():
        ph = ph_by_idx.get(idx)
        if ph is None:
            warnings.append(_no_idx(slide_index, idx, "text"))
            continue
        _set_text(ph, str(value))

    # 3. idx-keyed bullets — set placeholder <idx> to bullets (string or list).
    for idx, value in _idx_map(
        slide_spec.get("bullets"), "bullets", slide_index
    ).items():
        ph = ph_by_idx.get(idx)
        if ph is None:
            warnings.append(_no_idx(slide_index, idx, "bullets"))
            continue
        _set_body(ph, value)

    # 4. image — string (named placeholder) or dict directive (idx / box).
    if "image" in slide_spec and slide_spec["image"] is not None:
        _handle_image(
            slide, ph_by_idx, slide_spec["image"], base_dir, slide_index, warnings
        )

    # 5. force_font — set the run font of placeholder <idx> (brand fonts).
    _apply_force_font(ph_by_idx, slide_spec.get("force_font"), slide_index, warnings)


def _idx_map(value: Any, field: str, slide_index: int) -> dict[int, Any]:
    """Normalize an idx-keyed directive to ``{int: value}``."""
    if value is None:
        return {}
    if not isinstance(value, dict):
        raise BuildError(
            f"slide #{slide_index}: ``{field}`` must be a mapping of "
            f"placeholder-idx -> value"
        )
    out: dict[int, Any] = {}
    for key, val in value.items():
        try:
            out[int(key)] = val
        except (TypeError, ValueError):
            raise BuildError(
                f"slide #{slide_index}: ``{field}`` key {key!r} is not an "
                f"integer placeholder idx"
            ) from None
    return out


def _no_idx(slide_index: int, idx: int, field: str) -> str:
    return (
        f"slide #{slide_index}: {field} idx {idx} skipped (layout has no "
        f"placeholder with that idx — run deck-catalog)"
    )


def _handle_image(
    slide: Any,
    ph_by_idx: dict[int, Any],
    value: Any,
    base_dir: Path,
    slide_index: int,
    warnings: list[str],
) -> None:
    """Insert an image. ``value`` is either:

    - a **string** path -> routed to a named image/picture placeholder
      (backward-compatible simple form);
    - a **dict** ``{path, idx}`` -> inserted into placeholder ``idx``;
    - a **dict** ``{path, left, top, width}`` -> free picture at an explicit
      position in inches (no placeholder needed — for diagram slides).
    """
    if isinstance(value, str):
        image_path = _resolve_image(base_dir, value, slide_index)
        ph = _find_placeholder(slide, ("image", "picture")) or (
            _find_first_picture_placeholder(slide)
        )
        if ph is None:
            warnings.append(
                f"slide #{slide_index}: image given but no picture placeholder "
                f"found — use image: {{path, left, top, width}} to place it freely"
            )
            return
        _set_image(slide, ph, image_path)
        return

    if not isinstance(value, dict):
        raise BuildError(
            f"slide #{slide_index}: ``image`` must be a path string or a mapping"
        )

    path = value.get("path")
    if not path:
        raise BuildError(f"slide #{slide_index}: image.path is required")
    image_path = _resolve_image(base_dir, str(path), slide_index)

    if "idx" in value:
        ph = ph_by_idx.get(int(value["idx"]))
        if ph is None:
            warnings.append(_no_idx(slide_index, int(value["idx"]), "image"))
            return
        _set_image(slide, ph, image_path)
        return

    if all(k in value for k in ("left", "top", "width")):
        from pptx.util import Inches

        kwargs: dict[str, Any] = {"width": Inches(float(value["width"]))}
        if value.get("height") is not None:
            kwargs["height"] = Inches(float(value["height"]))
        slide.shapes.add_picture(
            str(image_path),
            Inches(float(value["left"])),
            Inches(float(value["top"])),
            **kwargs,
        )
        return

    raise BuildError(
        f"slide #{slide_index}: image dict needs either ``idx`` (a placeholder) "
        f"or ``left`` + ``top`` + ``width`` (an explicit position in inches)"
    )


def _resolve_image(base_dir: Path, rel_or_abs: str, slide_index: int) -> Path:
    image_path = (base_dir / rel_or_abs).resolve()
    if not image_path.exists():
        raise BuildError(f"slide #{slide_index}: image not found: {image_path}")
    return image_path


def _apply_force_font(
    ph_by_idx: dict[int, Any],
    entries: Any,
    slide_index: int,
    warnings: list[str],
) -> None:
    """Override the run font of placeholders by idx — for brand fonts the
    layout does not apply to populated text. Each entry is
    ``{idx: <int>, name: <font name>}``."""
    if not entries:
        return
    if not isinstance(entries, list):
        raise BuildError(
            f"slide #{slide_index}: ``force_font`` must be a list of "
            f"{{idx, name}} entries"
        )
    for entry in entries:
        if not isinstance(entry, dict) or "idx" not in entry or "name" not in entry:
            raise BuildError(
                f"slide #{slide_index}: each ``force_font`` entry needs ``idx`` "
                f"and ``name``"
            )
        ph = ph_by_idx.get(int(entry["idx"]))
        if ph is None or not ph.has_text_frame:
            warnings.append(_no_idx(slide_index, int(entry["idx"]), "force_font"))
            continue
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = str(entry["name"])
